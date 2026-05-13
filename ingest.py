import os
import pickle
import pdfplumber
import openpyxl
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument
from pdf2image import convert_from_path
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings
from langchain_chroma import Chroma

load_dotenv()

DATA_PATH = "data/"
VECTORSTORE_PATH = "vectorstore/"
DOCUMENTS_PATH = "vectorstore/documents.pkl"

TABLE_ROW_CHUNK_SIZE = 50

# ── Save / load ───────────────────────────────────────────────────────────────

def save_documents(documents, path=DOCUMENTS_PATH):
    with open(path, "wb") as f:
        pickle.dump(documents, f)

def load_saved_documents(path=DOCUMENTS_PATH):
    with open(path, "rb") as f:
        return pickle.load(f)

# ── PDF ───────────────────────────────────────────────────────────────────────

def is_presentation_pdf(filepath):
    with pdfplumber.open(filepath) as pdf:
        if pdf.pages:
            ratio = pdf.pages[0].width / pdf.pages[0].height
            return ratio > 1.2
    return False

def load_pdf_with_pdfplumber(filepath):
    documents = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                documents.append(Document(
                    page_content=text.strip(),
                    metadata={"source": filepath, "page": i, "doc_type": "document"}
                ))
    return documents

def load_pdf_with_ocr(filepath):
    documents = []
    images = convert_from_path(filepath, dpi=150)
    for i, image in enumerate(images):
        text = pytesseract.image_to_string(image, lang="eng+deu")
        if text and text.strip():
            documents.append(Document(
                page_content=text.strip(),
                metadata={"source": filepath, "page": i, "doc_type": "presentation"}
            ))
    return documents

def load_pdf(filepath):
    if is_presentation_pdf(filepath):
        print(f"    → detected as presentation, using OCR")
        return load_pdf_with_ocr(filepath)
    else:
        print(f"    → detected as document, using pdfplumber")
        return load_pdf_with_pdfplumber(filepath)

# ── PPTX ──────────────────────────────────────────────────────────────────────

def load_pptx(filepath):
    prs = Presentation(filepath)
    documents = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            documents.append(Document(
                page_content="\n".join(slide_text),
                metadata={"source": filepath, "page": i, "doc_type": "presentation"}
            ))
    return documents

# ── DOCX ──────────────────────────────────────────────────────────────────────

def load_docx(filepath):
    doc = DocxDocument(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
            if row_text:
                paragraphs.append(row_text)
    if paragraphs:
        return [Document(
            page_content="\n".join(paragraphs),
            metadata={"source": filepath, "doc_type": "document"}
        )]
    return []

# ── CSV ───────────────────────────────────────────────────────────────────────

def load_csv(filepath):
    df = pd.read_csv(filepath, encoding="utf-8")
    documents = []

    if len(df) <= TABLE_ROW_CHUNK_SIZE:
        documents.append(Document(
            page_content=df.to_string(index=False),
            metadata={"source": filepath, "doc_type": "table"}
        ))
    else:
        for i in range(0, len(df), TABLE_ROW_CHUNK_SIZE):
            chunk_df = df.iloc[i:i + TABLE_ROW_CHUNK_SIZE]
            documents.append(Document(
                page_content=chunk_df.to_string(index=False),
                metadata={"source": filepath, "doc_type": "table", "rows": f"{i}-{i + TABLE_ROW_CHUNK_SIZE}"}
            ))
    return documents

# ── XLSX ──────────────────────────────────────────────────────────────────────

def load_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath)
    documents = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = []
        headers = None

        for i, row in enumerate(ws.iter_rows(values_only=True)):
            row_text = " | ".join([str(c) for c in row if c is not None])
            if not row_text.strip():
                continue
            if i == 0:
                headers = row_text  # save header row
                all_rows.append(row_text)
            else:
                all_rows.append(row_text)

        if not all_rows:
            continue

        # Small sheet — keep as one chunk
        if len(all_rows) <= TABLE_ROW_CHUNK_SIZE:
            documents.append(Document(
                page_content="\n".join(all_rows),
                metadata={"source": filepath, "sheet": sheet_name, "doc_type": "table"}
            ))
        else:
            # Large sheet — split by rows, repeat headers in each chunk
            data_rows = all_rows[1:]  # skip header
            for i in range(0, len(data_rows), TABLE_ROW_CHUNK_SIZE):
                chunk_rows = data_rows[i:i + TABLE_ROW_CHUNK_SIZE]
                chunk_text = (headers + "\n" if headers else "") + "\n".join(chunk_rows)
                documents.append(Document(
                    page_content=chunk_text,
                    metadata={"source": filepath, "sheet": sheet_name, "doc_type": "table", "rows": f"{i}-{i + TABLE_ROW_CHUNK_SIZE}"}
                ))

    return documents

# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_documents(documents):
    prose_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=150
    )
    chunks = []
    for doc in documents:
        doc_type = doc.metadata.get("doc_type", "document")
        if doc_type in ("presentation", "table"):
            chunks.append(doc)
        elif doc_type == "document":
            chunks.extend(prose_splitter.split_documents([doc]))
    return chunks

# ── Main loader ───────────────────────────────────────────────────────────────

def load_documents():
    all_documents = []
    file_count = 0

    for root, dirs, files in os.walk(DATA_PATH):
        for file in files:
            file_path = os.path.join(root, file)
            try:
                if file.endswith(".pdf"):
                    docs = load_pdf(file_path)
                elif file.endswith(".pptx"):
                    docs = load_pptx(file_path)
                elif file.endswith(".docx"):
                    docs = load_docx(file_path)
                elif file.endswith(".csv"):
                    docs = load_csv(file_path)
                elif file.endswith(".xlsx"):
                    docs = load_xlsx(file_path)
                elif file.endswith(".txt") or file.endswith(".md"):
                    loader = TextLoader(file_path, encoding="utf-8")
                    docs = loader.load()
                    for doc in docs:
                        doc.metadata["doc_type"] = "document"
                else:
                    print(f"  Skipping: {file}")
                    continue

                file_count += 1
                print(f"  {len(docs)} pages — {file}")
                all_documents.extend(docs)

            except Exception as e:
                print(f"  Error loading {file}: {e}")
        break
    print(f"\nLoaded {file_count} files ({len(all_documents)} pages total)")
    return all_documents

# ── Ingest ────────────────────────────────────────────────────────────────────

def ingest_documents():
    print("Loading documents...")
    documents = load_documents()

    print("\nChunking documents...")
    chunks = chunk_documents(documents)
    print(f"Created {len(chunks)} chunks")

    print("\nEmbedding and storing in vector database...")
    embeddings = OpenAIEmbeddings()
    vectorstore = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        collection_name="knowledge_base",
        persist_directory=VECTORSTORE_PATH
    )

    print("Saving documents for BM25...")
    save_documents(chunks)
    print("Done!")


def ingest_single_file(filepath):
    """Ingest a single file and add it to the existing vector store."""
    filename = os.path.basename(filepath)
    print(f"Ingesting {filename}...")
    
    # Load the file
    try:
        if filepath.endswith(".pdf"):
            docs = load_pdf(filepath)
        elif filepath.endswith(".pptx"):
            docs = load_pptx(filepath)
        elif filepath.endswith(".docx"):
            docs = load_docx(filepath)
        elif filepath.endswith(".csv"):
            docs = load_csv(filepath)
        elif filepath.endswith(".xlsx"):
            docs = load_xlsx(filepath)
        elif filepath.endswith(".txt") or filepath.endswith(".md"):
            loader = TextLoader(filepath, encoding="utf-8")
            docs = loader.load()
            for doc in docs:
                doc.metadata["doc_type"] = "document"
        else:
            raise ValueError(f"Unsupported file type: {filename}")
    except Exception as e:
        raise Exception(f"Could not load {filename}: {e}")

    # Chunk the new docs
    chunks = chunk_documents(docs)
    print(f"Created {len(chunks)} chunks from {filename}")

    # Add to existing vector store
    embeddings = OpenAIEmbeddings()
    vectorstore = Chroma(
        persist_directory=VECTORSTORE_PATH,
        embedding_function=embeddings,
        collection_name="knowledge_base"
    )
    vectorstore.add_documents(chunks)
    print(f"Added to vector store")

    # Update pickle file
    existing_docs = load_saved_documents()
    updated_docs = existing_docs + chunks
    save_documents(updated_docs)
    print(f"Updated BM25 documents — total: {len(updated_docs)} chunks")

    return len(chunks)
if __name__ == "__main__":
    ingest_documents()